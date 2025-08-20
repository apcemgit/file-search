"""
Handles reading content from various file types.
"""
import os
import csv
import pandas as pd

# Optional dependencies
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

def read_file_content(filepath):
    """Extract text from supported file types."""
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        elif ext == '.pdf':
            if PdfReader is None:
                return "[PDF support not installed. Please run: pip install PyPDF2]"
            with open(filepath, 'rb') as f:
                reader = PdfReader(f)
                return " ".join((page.extract_text() or "") for page in reader.pages)

        elif ext == '.docx':
            if Document is None:
                return "[DOCX support not installed. Please run: pip install python-docx]"
            doc = Document(filepath)
            return " ".join(p.text for p in doc.paragraphs)

        elif ext == '.pptx':
            if Presentation is None:
                return "[PPTX support not installed. Please run: pip install python-pptx]"
            prs = Presentation(filepath)
            return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

        elif ext == '.csv':
            text = ""
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    text += " ".join(row) + " "
            return text

        elif ext == '.xlsx':
            try:
                df = pd.read_excel(filepath, sheet_name=None)  # All sheets
                text = ""
                for sheet_name, sheet_df in df.items():
                    text += f"Sheet: {sheet_name} " + sheet_df.to_string() + " "
                return text
            except Exception as e:
                return f"[Excel read error: {str(e)}]"

        else:
            return f"[Unsupported format: {ext}]"
    except Exception as e:
        return f"[Error reading file: {str(e)}]"
