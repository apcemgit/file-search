import os
import argparse
from datetime import datetime

# Optional dependencies – check if installed
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def read_file_content(filepath):
    """
    Extract text from supported file types.
    Returns text content or None if unsupported/failed.
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        elif ext == '.pdf':
            if PdfReader is None:
                print(f"Skipped (PDF support not installed): {filepath}")
                return ""
            reader = PdfReader(filepath)
            text = ""
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + " "
            return text

        elif ext == '.docx':
            if Document is None:
                print(f"Skipped (DOCX support not installed): {filepath}")
                return ""
            doc = Document(filepath)
            return " ".join(paragraph.text for paragraph in doc.paragraphs)

        elif ext == '.pptx':
            if Presentation is None:
                print(f"Skipped (PPTX support not installed): {filepath}")
                return ""
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + " "
            return text

        else:
            return ""  # Unsupported format
    except Exception as e:
        print(f"Error reading {filepath}: {e}")
        return ""


def search_files(directory, keywords, extensions=None, match_all=True,
                 case_sensitive=False, search_content=False):
    """
    Search files by filename and/or content.
    Returns list of tuples: (filepath, size, mtime)
    """
    matched_files = []

    search_keywords = [kw.lower() for kw in keywords] if not case_sensitive else keywords

    for root, dirs, files in os.walk(directory):
        for file in files:
            filepath = os.path.join(root, file)
            filename = file if case_sensitive else file.lower()
            file_ext = os.path.splitext(file)[1][1:].lower()

            # Filter by extension
            if extensions and file_ext not in ([e.lower() for e in extensions] if not case_sensitive else extensions):
                continue

            # Check filename match
            name_matches = (
                all(k.lower() in filename for k in keywords) if match_all
                else any(k.lower() in filename for k in keywords)
            )

            content_matches = False
            if search_content:
                content = read_file_content(filepath)
                content_str = content if case_sensitive else content.lower()
                content_matches = (
                    all(k.lower() in content_str for k in keywords) if match_all
                    else any(k.lower() in content_str for k in keywords)
                )

            # Match if: (searching content and content matches) OR (not searching content and filename matches)
            if (search_content and content_matches) or (not search_content and name_matches):
                stat = os.stat(filepath)
                matched_files.append({
                    'path': filepath,
                    'size': stat.st_size,
                    'mtime': stat.st_mtime
                })

    return matched_files


def main():
    parser = argparse.ArgumentParser(description="Advanced file search by name/content with sorting.")
    parser.add_argument('keywords', nargs='+', help="Keywords to search for")
    parser.add_argument('-d', '--directory', default='.', help="Directory to search (default: current)")
    parser.add_argument('-e', '--ext', nargs='+', help="Filter by extensions: pdf docx txt pptx etc.")
    parser.add_argument('--match-any', action='store_true', help="Match any keyword (default: all)")
    parser.add_argument('--case-sensitive', action='store_true', help="Case-sensitive search")
    parser.add_argument('--content', action='store_true', help="Search inside file content (supports txt, pdf, docx, pptx)")
    parser.add_argument('--sort', choices=['name', 'date', 'size', 'none'], default='name',
                        help="Sort results by name, date (newest first), size (largest first), or none")

    args = parser.parse_args()

    # Validate directory
    if not os.path.isdir(args.directory):
        print(f"Error: Directory '{args.directory}' does not exist.")
        return

    print(f"🔍 Searching in: {os.path.abspath(args.directory)}")
    print(f"📌 Keywords: {args.keywords}")
    print(f"📎 Extensions: {args.ext if args.ext else 'All'}")
    print(f"🔤 Match mode: {'Any keyword' if args.match_any else 'All keywords'}")
    print(f"🔠 Case sensitive: {args.case_sensitive}")
    print(f"📄 Search in content: {args.content}")
    print(f"📊 Sort by: {args.sort.capitalize()}")
    print("-" * 60)

    # Perform search
    results = search_files(
        directory=args.directory,
        keywords=args.keywords,
        extensions=args.ext,
        match_all=not args.match_any,
        case_sensitive=args.case_sensitive,
        search_content=args.content
    )

    # Sort results
    if args.sort == 'name':
        results.sort(key=lambda x: os.path.basename(x['path']).lower())
    elif args.sort == 'date':
        results.sort(key=lambda x: x['mtime'], reverse=True)  # Newest first
    elif args.sort == 'size':
        results.sort(key=lambda x: x['size'], reverse=True)  # Largest first
    # 'none' means no sort

    # Display results
    if results:
        print("✅ Found matching files:")
        for item in results:
            basename = os.path.basename(item['path'])
            size_kb = item['size'] // 1024
            date_str = datetime.fromtimestamp(item['mtime']).strftime("%Y-%m-%d %H:%M")
            print(f"  [{size_kb:6d} KB | {date_str}] {item['path']}")
    else:
        print("❌ No matching files found.")


if __name__ == "__main__":
    main()