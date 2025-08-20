import os
import re
import sys
import csv
import webbrowser
import pandas as pd
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, font
from datetime import datetime


# Optional dependencies
try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class Tooltip:
    def __init__(self, widget, text):
        self.widget = widget
        self.text = text
        self.tooltip_window = None
        self.widget.bind("<Enter>", self.show_tooltip)
        self.widget.bind("<Leave>", self.hide_tooltip)

    def show_tooltip(self, event=None):
        if self.tooltip_window:
            return
        x, y, _, _ = self.widget.bbox("insert") if hasattr(self.widget, 'bbox') else (0, 0, 0, 0)
        x += self.widget.winfo_rootx() + 25
        y += self.widget.winfo_rooty() + 25

        self.tooltip_window = tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        label = tk.Label(tw, text=self.text, justify="left",
                         background="#ffffe0", relief="solid", borderwidth=1,
                         font=("Tahoma", "9"), padx=5, pady=3)
        label.pack()

    def hide_tooltip(self, event=None):
        if self.tooltip_window:
            self.tooltip_window.destroy()
            self.tooltip_window = None


class FileSearchApp:
    def __init__(self, root):
        self.root = root
        self.root.title("📁 Advanced File Search Pro")
        self.root.geometry("1100x800")
        self.root.configure(padx=15, pady=15)

        # Fonts
        self.default_font = font.nametofont("TkDefaultFont")
        self.bold_font = font.Font(**self.default_font.config())
        self.mono_font = font.Font(family="Courier", size=10)

        # Icons (emoji-style fallback)
        self.icons = {
            'txt': '📄',
            'pdf': '📕',
            'docx': '📘',
            'pptx': '📊',
            'xlsx': '📈',
            'csv': '📑',
            'py': '🐍',
            'js': '🟨',
            'json': '🔧',
            'md': '📝',
            'jpg': '🖼️',
            'png': '🖼️',
            'default': '📎'
        }

        self.results = []
        self.create_widgets()

    def create_widgets(self):
        # === Title ===
        title = ttk.Label(self.root, text="Advanced File Search Pro", font=("Helvetica", 16, "bold"))
        title.grid(row=0, column=0, columnspan=5, pady=(0, 15))

        # === Directory Selection ===
        ttk.Label(self.root, text="Search Directory:").grid(row=1, column=0, sticky="w", pady=5)
        self.directory_var = tk.StringVar(value=os.getcwd())
        dir_entry = ttk.Entry(self.root, textvariable=self.directory_var, width=50)
        dir_entry.grid(row=1, column=1, columnspan=3, padx=5, pady=5, sticky="ew")
        ttk.Button(self.root, text="Browse...", command=self.browse_directory).grid(row=1, column=4, padx=5)

        # === Keywords ===
        ttk.Label(self.root, text="Pattern:").grid(row=2, column=0, sticky="w", pady=5)
        self.pattern_var = tk.StringVar()
        pattern_entry = ttk.Entry(self.root, textvariable=self.pattern_var, width=50)
        pattern_entry.grid(row=2, column=1, columnspan=3, padx=5, pady=5, sticky="ew")
        Tooltip(pattern_entry, "Enter text or regex pattern\nUse space-separated words (non-regex)\nor a full regex like: report_.*2024\\.pdf")

        # === Extensions ===
        ttk.Label(self.root, text="Extensions:").grid(row=3, column=0, sticky="w", pady=5)
        self.ext_var = tk.StringVar(value="txt pdf docx pptx xlsx csv py js")
        ext_entry = ttk.Entry(self.root, textvariable=self.ext_var, width=50)
        ext_entry.grid(row=3, column=1, columnspan=3, padx=5, pady=5, sticky="ew")
        Tooltip(ext_entry, "Filter by extensions: pdf docx txt pptx xlsx csv py js\nLeave empty to include all")

        # === Options Frame ===
        options_frame = ttk.LabelFrame(self.root, text="Search Options", padding=10)
        options_frame.grid(row=4, column=0, columnspan=5, pady=10, sticky="ew")

        self.match_any_var = tk.BooleanVar()
        self.case_sensitive_var = tk.BooleanVar()
        self.search_content_var = tk.BooleanVar()
        self.regex_var = tk.BooleanVar()
        self.sort_var = tk.StringVar(value="name")

        ttk.Checkbutton(options_frame, text="Match any keyword", variable=self.match_any_var).grid(row=0, column=0, padx=10)
        ttk.Checkbutton(options_frame, text="Case sensitive", variable=self.case_sensitive_var).grid(row=0, column=1, padx=10)
        ttk.Checkbutton(options_frame, text="Search in content", variable=self.search_content_var).grid(row=0, column=2, padx=10)
        ttk.Checkbutton(options_frame, text="Use Regex", variable=self.regex_var).grid(row=0, column=3, padx=10)

        ttk.Label(options_frame, text="Sort by:").grid(row=0, column=4, padx=5)
        ttk.Combobox(options_frame, textvariable=self.sort_var, values=["name", "date", "size", "none"],
                     state="readonly", width=10).grid(row=0, column=5, padx=5)

        # === Buttons & Progress ===
        btn_frame = ttk.Frame(self.root)
        btn_frame.grid(row=5, column=0, columnspan=5, pady=10, sticky="ew")

        self.search_btn = ttk.Button(btn_frame, text="🔍 Start Search", command=self.start_search)
        self.search_btn.pack(side="left", padx=2)

        self.export_btn = ttk.Button(btn_frame, text="💾 Export Results", command=self.export_results)
        self.export_btn.pack(side="left", padx=2)
        self.export_btn.config(state="disabled")

        ttk.Label(btn_frame, text="Scanned:").pack(side="right")
        self.scanned_label = ttk.Label(btn_frame, text="0")
        self.scanned_label.pack(side="right", padx=(0, 5))

        self.progress_var = tk.DoubleVar()
        self.progress_bar = ttk.Progressbar(btn_frame, variable=self.progress_var, maximum=100, mode="determinate")
        self.progress_bar.pack(side="right", fill="x", expand=True, padx=(5, 10))

        # === Results & Preview Split ===
        paned = ttk.PanedWindow(self.root, orient=tk.HORIZONTAL)
        paned.grid(row=6, column=0, columnspan=5, pady=10, sticky="nsew")
        self.root.grid_rowconfigure(6, weight=1)
        self.root.grid_columnconfigure(4, weight=1)

        # --- Left: Results Tree ---
        left_frame = ttk.Frame(paned)
        columns = ("Icon", "Name", "Size", "Modified")
        self.tree = ttk.Treeview(left_frame, columns=columns, show="headings", height=20)
        self.tree.heading("Icon", text="")
        self.tree.heading("Name", text="Name")
        self.tree.heading("Size", text="Size (KB)")
        self.tree.heading("Modified", text="Modified")
        self.tree.column("Icon", width=30, anchor="center")
        self.tree.column("Name", width=250)
        self.tree.column("Size", width=100, anchor="center")
        self.tree.column("Modified", width=150, anchor="center")

        vsb1 = ttk.Scrollbar(left_frame, orient="vertical", command=self.tree.yview)
        self.tree.configure(yscrollcommand=vsb1.set)
        self.tree.pack(side="left", fill="both", expand=True)
        vsb1.pack(side="right", fill="y")

        self.tree.bind("<Double-1>", self.open_selected_file)
        self.tree.bind("<<TreeviewSelect>>", self.show_preview)
        self.tree.bind("<Button-3>", self.show_context_menu)
        if sys.platform == "darwin":
            self.tree.bind("<Button-2>", self.show_context_menu)

        # --- Right: Preview Pane ---
        right_frame = ttk.Frame(paned)
        ttk.Label(right_frame, text="📄 Content Preview", font=("Helvetica", 11, "bold")).pack(anchor="w", pady=5)
        self.preview_text = tk.Text(right_frame, wrap="word", font=self.mono_font, height=20, bg="#f4f4f4")
        self.preview_text.pack(fill="both", expand=True, padx=5, pady=5)
        vsb2 = ttk.Scrollbar(self.preview_text, command=self.preview_text.yview)
        self.preview_text.configure(yscrollcommand=vsb2.set)
        vsb2.pack(side="right", fill="y")

        paned.add(left_frame, weight=3)
        paned.add(right_frame, weight=2)

        # === Footer ===
        footer_frame = ttk.Frame(self.root)
        footer_frame.grid(row=7, column=0, columnspan=5, pady=10)

        license_label = ttk.Label(footer_frame, text="License: MIT | Made with love: ", foreground="gray")
        license_label.pack(side="left")

        creator_label = tk.Label(footer_frame, text="Jhenbert", foreground="blue", cursor="hand2", font=("Helvetica", 9, "underline"))
        creator_label.pack(side="left")
        creator_label.bind("<Button-1>", lambda e: webbrowser.open("https://jhenbert.dev"))

        copyright_label = ttk.Label(footer_frame, text=f" © {datetime.now().year}", foreground="gray")
        copyright_label.pack(side="left")

        # Create context menu
        self.create_context_menu()

    def create_context_menu(self):
        self.context_menu = tk.Menu(self.root, tearoff=0)
        self.context_menu.add_command(label="📂 Open File", command=self.open_selected_file)
        self.context_menu.add_command(label="📁 Open File Location", command=self.open_file_location)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="📋 Copy Path", command=self.copy_path_to_clipboard)

    def show_context_menu(self, event):
        item = self.tree.identify_row(event.y)
        if item:
            if not self.tree.selection() or item not in self.tree.selection():
                self.tree.selection_set(item)
        else:
            return
        try:
            self.context_menu.tk_popup(event.x_root, event.y_root)
        finally:
            self.context_menu.grab_release()

    def browse_directory(self):
        path = filedialog.askdirectory(initialdir=self.directory_var.get())
        if path:
            self.directory_var.set(path)

    def read_file_content(self, filepath):
        """Extract text from supported file types including .csv and .xlsx"""
        ext = os.path.splitext(filepath)[1].lower()

        try:
            if ext == '.txt':
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    return f.read()

            elif ext == '.pdf':
                if PdfReader is None:
                    return "[PDF support not installed]"
                reader = PdfReader(filepath)
                return " ".join((page.extract_text() or "") for page in reader.pages)

            elif ext == '.docx':
                if Document is None:
                    return "[DOCX support not installed]"
                doc = Document(filepath)
                return " ".join(p.text for p in doc.paragraphs)

            elif ext == '.pptx':
                if Presentation is None:
                    return "[PPTX support not installed]"
                prs = Presentation(filepath)
                return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

            elif ext == '.csv':
                text = ""
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    for row in reader:
                        text += " ".join(row) + " "
                return text

            elif ext == '.xlsx':
                try:
                    df = pd.read_excel(filepath, sheet_name=None)  # All sheets
                    text = ""
                    for sheet_name, sheet_df in df.items():
                        text += f"Sheet: {sheet_name} " + sheet_df.to_string() + " "
                    return text
                except Exception as e:
                    return f"[Excel read error: {str(e)}]"

            else:
                return f"[Unsupported format: {ext}]"
        except Exception as e:
            return f"[Error reading file: {str(e)}]"

    def matches_pattern(self, text, pattern, use_regex=False, case_sensitive=False):
        if use_regex:
            flags = 0 if case_sensitive else re.IGNORECASE
            try:
                return re.search(pattern, text, flags) is not None
            except re.error as e:
                messagebox.showwarning("Regex Error", f"Invalid regex pattern:\n{str(e)}")
                return False
        else:
            keywords = pattern.split()
            if not case_sensitive:
                text = text.lower()
                keywords = [k.lower() for k in keywords]
            return all(k in text for k in keywords) if not self.match_any_var.get() \
                else any(k in text for k in keywords)

    def highlight_text(self, text, pattern, use_regex=False, case_sensitive=False):
        if use_regex:
            flags = 0 if case_sensitive else re.IGNORECASE
            try:
                return re.sub(f"({pattern})", r"**\1**", text, flags=flags)
            except re.error:
                return text
        else:
            for kw in pattern.split():
                if not case_sensitive:
                    text = re.sub(re.escape(kw), f"**{kw}**", text, flags=re.IGNORECASE)
                else:
                    text = text.replace(kw, f"**{kw}**")
            return text

    def start_search(self):
        for item in self.tree.get_children():
            self.tree.delete(item)
        self.results = []
        self.preview_text.delete(1.0, tk.END)

        directory = self.directory_var.get()
        pattern = self.pattern_var.get().strip()
        extensions = self.ext_var.get().strip().split() if self.ext_var.get().strip() else None
        match_any = self.match_any_var.get()
        case_sensitive = self.case_sensitive_var.get()
        search_content = self.search_content_var.get()
        use_regex = self.regex_var.get()
        sort_by = self.sort_var.get()

        if not pattern:
            messagebox.showwarning("Input Error", "Please enter a search pattern.")
            return

        if not os.path.isdir(directory):
            messagebox.showerror("Directory Error", f"Directory not found:\n{directory}")
            return

        self.search_btn.config(state="disabled")
        self.progress_var.set(0)
        self.scanned_label.config(text="0")
        self.root.update()

        total_files = sum(len(files) for _, _, files in os.walk(directory))
        if total_files == 0:
            messagebox.showinfo("No Files", "No files found in selected directory.")
            self.search_btn.config(state="normal")
            return

        scanned = 0
        for root_dir, dirs, files in os.walk(directory):
            for file in files:
                filepath = os.path.join(root_dir, file)
                filename = file if case_sensitive else file.lower()
                file_ext = os.path.splitext(file)[1][1:].lower()

                if extensions and file_ext not in ([e.lower() for e in extensions] if not case_sensitive else extensions):
                    scanned += 1
                    continue

                name_match = self.matches_pattern(filename, pattern, use_regex, case_sensitive)

                content_match = False
                content_snippet = ""
                if search_content:
                    content = self.read_file_content(filepath)
                    content_for_search = content if case_sensitive else content.lower()
                    pattern_for_search = pattern if case_sensitive or use_regex else pattern.lower()
                    content_match = self.matches_pattern(content_for_search, pattern_for_search, use_regex, case_sensitive)
                    if content_match:
                        pos = max(0, content.lower().find(pattern.lower()) - 50)
                        end = min(len(content), pos + 150)
                        content_snippet = content[pos:end].strip()

                if (search_content and content_match) or (not search_content and name_match):
                    try:
                        stat = os.stat(filepath)
                        result = {
                            'name': file,
                            'path': filepath,
                            'size': stat.st_size,
                            'mtime': stat.st_mtime,
                            'snippet': content_snippet,
                            'ext': file_ext
                        }
                        self.results.append(result)
                    except:
                        pass

                scanned += 1
                self.progress_var.set((scanned / total_files) * 100)
                self.scanned_label.config(text=str(scanned))
                if scanned % 20 == 0:
                    self.root.update_idletasks()

        if sort_by == "name":
            self.results.sort(key=lambda x: x['name'].lower())
        elif sort_by == "date":
            self.results.sort(key=lambda x: x['mtime'], reverse=True)
        elif sort_by == "size":
            self.results.sort(key=lambda x: x['size'], reverse=True)

        for r in self.results:
            size_kb = r['size'] // 1024
            date_str = datetime.fromtimestamp(r['mtime']).strftime("%Y-%m-%d %H:%M")
            icon = self.icons.get(r['ext'], self.icons['default'])
            highlighted_name = self.highlight_text(r['name'], pattern, use_regex, case_sensitive)
            self.tree.insert("", "end", values=(icon, highlighted_name, size_kb, date_str), tags=(r['path'],))

        self.progress_var.set(100)
        self.export_btn.config(state="normal")
        self.search_btn.config(state="normal")
        messagebox.showinfo("Done", f"Found {len(self.results)} file(s).")

    def open_selected_file(self, event=None):
        selected = self.tree.selection()
        if not selected:
            return
        item = self.tree.item(selected[0])
        if 'tags' not in item or not item['tags']:
            messagebox.showwarning("Error", "File path not available.")
            return
        filepath = item['tags'][0]
        if not os.path.exists(filepath):
            messagebox.showerror("Not Found", f"File does not exist:\n{filepath}")
            return
        try:
            if sys.platform == "win32":
                os.startfile(filepath)
            elif sys.platform == "darwin":
                os.system(f"open '{filepath}'")
            else:
                os.system(f"xdg-open '{filepath}'")
        except Exception as e:
            messagebox.showerror("Open Failed", f"Could not open file:\n{str(e)}")

    def open_file_location(self):
        selected = self.tree.selection()
        if not selected:
            return
        item = self.tree.item(selected[0])
        if 'tags' not in item or not item['tags']:
            messagebox.showwarning("Error", "File path not available.")
            return
        filepath = item['tags'][0]
        folder_path = os.path.dirname(filepath)
        if not os.path.exists(folder_path):
            messagebox.showerror("Not Found", f"Folder does not exist:\n{folder_path}")
            return
        try:
            if sys.platform == "win32":
                os.startfile(folder_path)
            elif sys.platform == "darwin":
                os.system(f"open '{folder_path}'")
            else:
                os.system(f"xdg-open '{folder_path}'")
        except Exception as e:
            messagebox.showerror("Open Failed", f"Could not open folder:\n{str(e)}")

    def copy_path_to_clipboard(self):
        selected = self.tree.selection()
        if not selected:
            return
        item = self.tree.item(selected[0])
        if 'tags' not in item or not item['tags']:
            return
        filepath = item['tags'][0]
        self.root.clipboard_clear()
        self.root.clipboard_append(filepath)
        self.root.update()
        messagebox.showinfo("Copied", "File path copied to clipboard!")

    def show_preview(self, event):
        self.preview_text.delete(1.0, tk.END)
        selected = self.tree.selection()
        if not selected:
            return
        item = self.tree.item(selected[0])
        filepath = item['tags'][0]
        result = next((r for r in self.results if r['path'] == filepath), None)

        if not result:
            self.preview_text.insert(tk.END, "⚠️ No data available.")
            return

        self.preview_text.insert(tk.END, f"📁 File: {os.path.basename(filepath)}\n")
        self.preview_text.insert(tk.END, f"🔗 Path: {filepath}\n\n")

        if not self.search_content_var.get():
            self.preview_text.insert(tk.END, "ℹ️ Content search was disabled. Re-run with 'Search in content' enabled.")
            return

        content = self.read_file_content(filepath)
        if len(content) < 2 or "support not installed" in content or "Error reading" in content:
            self.preview_text.insert(tk.END, f"❌ Could not read content:\n{content}")
            return

        highlighted = self.highlight_text(content, self.pattern_var.get().strip(), self.regex_var.get(), self.case_sensitive_var.get())
        try:
            match_pos = content.lower().find(self.pattern_var.get().strip().lower())
            start = max(0, match_pos - 100)
            end = min(len(content), match_pos + 300)
            snippet = highlighted[start:end]
            self.preview_text.insert(tk.END, "📌 Context around first match:\n\n")
            self.preview_text.insert(tk.END, snippet)
        except:
            self.preview_text.insert(tk.END, "🔤 Could not extract matching context.")

    def export_results(self):
        filepath = filedialog.asksaveasfilename(
            defaultextension=".xlsx",
            filetypes=[("Excel files", "*.xlsx"), ("CSV files", "*.csv"), ("All files", "*.*")]
        )
        if not filepath:
            return
        try:
            if filepath.endswith(".csv"):
                with open(filepath, 'w', encoding='utf-8', newline='') as f:
                    writer = csv.writer(f)
                    writer.writerow(["Filename", "Path", "Size (bytes)", "Modified", "Snippet"])
                    for r in self.results:
                        writer.writerow([
                            r['name'],
                            r['path'],
                            r['size'],
                            datetime.fromtimestamp(r['mtime']).strftime("%Y-%m-%d %H:%M:%S"),
                            r['snippet'][:500]
                        ])
            else:
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws.title = "Search Results"
                ws.append(["Filename", "Path", "Size (bytes)", "Modified", "Snippet"])
                for r in self.results:
                    ws.append([
                        r['name'],
                        r['path'],
                        r['size'],
                        datetime.fromtimestamp(r['mtime']).strftime("%Y-%m-%d %H:%M:%S"),
                        r['snippet'][:500]
                    ])
                wb.save(filepath)
            messagebox.showinfo("Export Success", f"Results saved to:\n{filepath}")
        except Exception as e:
            messagebox.showerror("Export Failed", f"Could not save file:\n{str(e)}")


if __name__ == "__main__":
    root = tk.Tk()
    app = FileSearchApp(root)
    root.mainloop()